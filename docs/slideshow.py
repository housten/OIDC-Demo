from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ---------- Helper Functions ----------
def add_slide(prs, title, bullets, notes, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
    slide.notes_slide.notes_text_frame.text = notes
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = RGBColor(255, 255, 255)
    return slide

# ---------- Create Presentation ----------
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "OIDC with GitHub Actions: Beta Introduction (Session 1)"
subtitle = slide.placeholders[1]
subtitle.text = "Presenter: Heidi Housten\nSolidify – now a part of Eficode"
for shape in slide.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255,255,255)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(24, 26, 27)

# ---------- Slides ----------
slides_content = [
("Agenda", [
    "Part 1: The Why and What (15 min)",
    "• Static Keys Problem",
    "• OIDC Solution & GitHub’s Role",
    "Part 2: Workload Identity & Trust (20 min)",
    "• JWT Deep Dive",
    "• Cloud Federation Demos",
    "Feedback & Q&A"
], "Outline of today's beta session for developers/devops."),

("The Hidden Risk: Static Credentials", [
    "AWS_ACCESS_KEY_ID / AZURE_SECRET",
    "Long-lived credentials; high risk if leaked",
    "Zero context on usage or origin"
], "Discuss static keys as a major CI/CD risk."),

("The OIDC Difference: Short-Lived Tokens", [
    "OIDC builds on OAuth 2.0",
    "Tokens issued by GitHub IdP, short-lived & signed",
    "No stored secrets — fresh token per run"
], "Explain OIDC temporary tokens and verifiable identity."),

("OIDC vs OAuth: Authentication vs Authorization", [
    "OAuth 2.0 → Authorization",
    "OIDC → Authentication",
    "CI/CD uses OIDC for identity → leads to cloud authorization"
], "Contrast the two protocols for clarity."),

("The GitHub Actions Identity Provider", [
    "Issuer: https://token.actions.githubusercontent.com",
    "permissions: id-token: write",
    "Use cloud login actions to exchange token"
], "GitHub acts as IdP issuing the JWT to cloud providers."),

("Workload Identity: Identity for Your Code", [
    "Human vs Machine identity",
    "Federation: trusting external IdP (GitHub)",
    "Azure Service Principal / AWS IAM Role"
], "Introduce workload identity concept."),

("Demo 1: The Well-Known Endpoint", [
    "URL: https://token.actions.githubusercontent.com/.well-known/openid-configuration",
    "Lists claims and public keys",
    "Proves GitHub is OIDC-compliant"
], "Show OIDC discovery metadata."),

("Demo 2: Inside the Identity Token", [
    "JWT = Header.Payload.Signature",
    "iss: GitHub | sub: repo/branch | aud: cloud service",
    "Use jwt.ms to inspect token"
], "Visualize token claims and meaning."),

("Demo 3: Configuring Azure Trust", [
    "Entra ID App Registration (Service Principal)",
    "Federated Credential links GitHub JWT to Entra ID",
    "Subject match creates trust bridge"
], "Walk through Entra ID setup for federation."),

("Feedback & Q&A", [
    "Your input shapes Session 2!",
    "Use Q&A or QR form for comments",
    "Thank you for participating in the beta!"
], "Invite audience feedback and wrap-up."),
]

for title, bullets, notes in slides_content:
    s = add_slide(prs, title, bullets, notes)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(24,26,27)

# ---------- Save ----------
prs.save("OIDC_Beta_Intro.pptx")
print("Saved as OIDC_Beta_Intro.pptx")
